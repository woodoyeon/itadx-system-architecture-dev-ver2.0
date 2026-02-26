# -*- coding: utf-8 -*-
"""ItaDX 프론트↔백 통신 가이드 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x1E, 0x29, 0x3B)      # 진한 남색 배경
BG_CARD = RGBColor(0x27, 0x34, 0x4A)      # 카드 배경
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)  # 밝은 파랑
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80) # 초록
ACCENT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)# 노랑
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)# 보라
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)# 오렌지
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)   # 빨강
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
VERY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
CODE_BG = RGBColor(0x0F, 0x17, 0x2A)      # 코드 블록 배경

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 둥근 정도
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_code_block(slide, left, top, width, height, lines, font_size=11):
    """코드 블록: 어두운 배경 + 모노스페이스 폰트"""
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x33, 0x44, 0x55), border_width=Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    # 기존 paragraph 삭제
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_info[0] if isinstance(line_info, tuple) else line_info
        clr = line_info[1] if isinstance(line_info, tuple) and len(line_info) > 1 else VERY_LIGHT
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(0)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(2)):
    """화살표 커넥터"""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    # 화살표 머리 (end)
    connector.end_x = x2
    connector.end_y = y2
    return connector


def add_circle(slide, left, top, size, fill_color, text="", font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    return shape


def add_pill(slide, left, top, width, height, fill_color, text, font_size=11, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# 제목
add_text(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.2),
         "ItaDX  프론트 ↔ 백엔드 통신 가이드", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.8),
         "이 문서만 보면  프론트·백엔드  개발자 모두\n같은 틀 안에서 바로 개발할 수 있습니다", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 하단 태그들
tags = [("REST API", ACCENT_BLUE), ("WebSocket", ACCENT_GREEN), ("JSON 규격", ACCENT_YELLOW), ("코드 예시", ACCENT_PURPLE)]
start_x = Inches(3.2)
for i, (tag, clr) in enumerate(tags):
    add_pill(slide, start_x + Inches(i * 1.9), Inches(4.5), Inches(1.6), Inches(0.4), clr, tag, font_size=13, text_color=BG_DARK)

add_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
         "ItaDX MVP  ·  2026.02", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 2: 전체 아키텍처 한눈에 보기
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "전체 아키텍처  —  요청이 어떻게 흘러가나?", font_size=28, bold=True, color=WHITE)
add_text(slide, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
         "프론트는 /api 만 부르면 됩니다. 나머지는 Next→Gateway→백엔드가 알아서 처리합니다.", font_size=14, color=LIGHT_GRAY)

# 박스들 그리기
box_y = Inches(1.7)
box_h = Inches(4.8)

# 1) 브라우저
add_shape(slide, Inches(0.4), box_y, Inches(2.2), box_h, fill_color=RGBColor(0x1E, 0x40, 0x5F), border_color=ACCENT_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.5), Inches(1.8), Inches(2.0), Inches(0.4), "🖥  브라우저", font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.3), Inches(2.0), Inches(1.2),
         "api.post('/auth/login')\napi.get('/receivings')\napi.patch('/receivings/:id/confirm')", font_size=11, color=VERY_LIGHT, font_name="Consolas", alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.6), Inches(2.0), Inches(0.3), "항상  /api  로 시작", font_size=12, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)
# JWT 자동
add_pill(slide, Inches(0.7), Inches(4.1), Inches(1.8), Inches(0.35), RGBColor(0x33, 0x44, 0x55), "JWT 자동 첨부", font_size=10, text_color=ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(4.55), Inches(2.0), Inches(0.3), "📂 lib/api.ts", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
# WebSocket 수신
add_pill(slide, Inches(0.7), Inches(5.1), Inches(1.8), Inches(0.35), RGBColor(0x1A, 0x3A, 0x2A), "WebSocket 수신", font_size=10, text_color=ACCENT_GREEN)
add_text(slide, Inches(0.5), Inches(5.55), Inches(2.0), Inches(0.5), "socket.on('receiving:confirmed')\n→ 캐시 무효화 → 화면 갱신", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 2) Next.js rewrite
add_shape(slide, Inches(3.2), box_y, Inches(2.0), Inches(2.0), fill_color=RGBColor(0x2D, 0x2D, 0x44), border_color=ACCENT_PURPLE, border_width=Pt(2))
add_text(slide, Inches(3.2), Inches(1.8), Inches(2.0), Inches(0.4), "⚡ Next.js", font_size=16, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(3.2), Inches(2.3), Inches(2.0), Inches(0.6), "/api/:path*\n   ↓ rewrite\nlocalhost:4003", font_size=11, color=VERY_LIGHT, font_name="Consolas", alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(3.2), Inches(3.1), Inches(2.0), Inches(0.3), "📂 next.config.js", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 3) Gateway
add_shape(slide, Inches(5.8), box_y, Inches(2.3), box_h, fill_color=RGBColor(0x3B, 0x2D, 0x1E), border_color=ACCENT_ORANGE, border_width=Pt(2))
add_text(slide, Inches(5.8), Inches(1.8), Inches(2.3), Inches(0.4), "🔀  Gateway :4003", font_size=16, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.8), Inches(2.3), Inches(2.3), Inches(0.4), "URL prefix 보고 분배", font_size=12, color=VERY_LIGHT, alignment=PP_ALIGN.CENTER)

routes_text = "/api/auth     → :4001\n/api/marts    → :4000\n/api/users    → :4000\n/api/receivings → :4002\n/api/settlements → :4002\n/api/v41      → :8000"
add_text(slide, Inches(5.9), Inches(2.8), Inches(2.1), Inches(2.0), routes_text, font_size=10, color=VERY_LIGHT, font_name="Consolas")

add_pill(slide, Inches(6.1), Inches(5.0), Inches(1.9), Inches(0.35), RGBColor(0x44, 0x33, 0x22), "Rate Limit 100/min", font_size=9, text_color=ACCENT_ORANGE)
add_text(slide, Inches(5.8), Inches(5.5), Inches(2.3), Inches(0.3), "📂 gateway-api/main.ts", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 4) 백엔드 서버들
add_shape(slide, Inches(8.7), box_y, Inches(4.2), box_h, fill_color=RGBColor(0x1E, 0x33, 0x1E), border_color=ACCENT_GREEN, border_width=Pt(2))
add_text(slide, Inches(8.7), Inches(1.8), Inches(4.2), Inches(0.4), "🏗  백엔드 서비스", font_size=16, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# auth-api 박스
add_shape(slide, Inches(8.9), Inches(2.4), Inches(1.9), Inches(1.5), fill_color=RGBColor(0x15, 0x28, 0x15), border_color=RGBColor(0x4A, 0x80, 0x4A), border_width=Pt(1))
add_text(slide, Inches(8.9), Inches(2.45), Inches(1.9), Inches(0.3), "auth-api :4001", font_size=12, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.9), Inches(2.8), Inches(1.9), Inches(0.9), "POST /login\nPOST /refresh\nPOST /logout\nGET  /profile", font_size=10, color=VERY_LIGHT, font_name="Consolas", alignment=PP_ALIGN.CENTER)

# admin-api 박스
add_shape(slide, Inches(11.0), Inches(2.4), Inches(1.8), Inches(1.5), fill_color=RGBColor(0x15, 0x28, 0x15), border_color=RGBColor(0x4A, 0x80, 0x4A), border_width=Pt(1))
add_text(slide, Inches(11.0), Inches(2.45), Inches(1.8), Inches(0.3), "admin-api :4000", font_size=12, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(11.0), Inches(2.8), Inches(1.8), Inches(0.9), "마트 CRUD\n지점 관리\n가맹점 관리\n대시보드", font_size=10, color=VERY_LIGHT, alignment=PP_ALIGN.CENTER)

# erp-api 박스
add_shape(slide, Inches(8.9), Inches(4.2), Inches(1.9), Inches(1.8), fill_color=RGBColor(0x15, 0x28, 0x15), border_color=ACCENT_YELLOW, border_width=Pt(1))
add_text(slide, Inches(8.9), Inches(4.25), Inches(1.9), Inches(0.3), "⭐ erp-api :4002", font_size=12, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.9), Inches(4.6), Inches(1.9), Inches(1.2), "★ 입고확인\n  정산 관리\n→ Bull Queue\n→ WebSocket 알림", font_size=10, color=VERY_LIGHT, alignment=PP_ALIGN.CENTER)

# engine-api 박스
add_shape(slide, Inches(11.0), Inches(4.2), Inches(1.8), Inches(1.8), fill_color=RGBColor(0x15, 0x28, 0x15), border_color=RGBColor(0x4A, 0x80, 0x4A), border_width=Pt(1))
add_text(slide, Inches(11.0), Inches(4.25), Inches(1.8), Inches(0.3), "engine-api :8000", font_size=12, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(11.0), Inches(4.6), Inches(1.8), Inches(1.2), "Python FastAPI\nv41 심사\nv10 신용점수\n듀얼트랙\n점포리스크", font_size=10, color=VERY_LIGHT, alignment=PP_ALIGN.CENTER)

# 화살표 (텍스트로 표현)
add_text(slide, Inches(2.55), Inches(2.3), Inches(0.8), Inches(0.4), "→", font_size=28, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.1), Inches(2.3), Inches(0.8), Inches(0.4), "→", font_size=28, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.0), Inches(2.8), Inches(0.8), Inches(0.4), "→", font_size=28, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# WebSocket 역방향 화살표
add_text(slide, Inches(2.55), Inches(5.2), Inches(6.0), Inches(0.4), "← ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─  WebSocket  ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ←", font_size=12, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 3: REST API 통신 규칙
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "REST API 통신 규칙  —  이것만 알면 됩니다", font_size=28, bold=True, color=WHITE)

# 요청 섹션
add_text(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.4),
         "📤  프론트 → 백  요청 (Request)", font_size=20, bold=True, color=ACCENT_BLUE)

add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), fill_color=RGBColor(0x15, 0x25, 0x3B), border_color=ACCENT_BLUE, border_width=Pt(1))

# 요청 구성요소
items = [
    ("① 메서드 + URL", "GET /api/receivings?page=1&status=pending", ACCENT_BLUE),
    ("② 헤더 (자동)", "Authorization: Bearer eyJhbG...\nContent-Type: application/json", ACCENT_PURPLE),
    ("③ 바디 (POST/PATCH만)", '{ "email": "admin@itadx.com",\n  "password": "password123" }', ACCENT_YELLOW),
]

y_pos = Inches(1.75)
for title, code, clr in items:
    add_text(slide, Inches(0.8), y_pos, Inches(5.4), Inches(0.35), title, font_size=14, bold=True, color=clr)
    y_pos += Inches(0.35)
    code_lines = [(line, VERY_LIGHT) for line in code.split("\n")]
    add_code_block(slide, Inches(0.8), y_pos, Inches(5.4), Inches(0.2 + len(code_lines) * 0.25), code_lines, font_size=11)
    y_pos += Inches(0.3 + len(code_lines) * 0.28)

# 프론트 코드 예시
add_text(slide, Inches(0.8), y_pos, Inches(5.4), Inches(0.35), "④ 프론트 코드 (이것만 쓰면 됨)", font_size=14, bold=True, color=ACCENT_GREEN)
y_pos += Inches(0.35)
add_code_block(slide, Inches(0.8), y_pos, Inches(5.4), Inches(1.3), [
    ("// lib/api.ts 의 axios 인스턴스 사용", LIGHT_GRAY),
    ("import api from '@/lib/api';", ACCENT_PURPLE),
    ("", WHITE),
    ("// GET (조회)", LIGHT_GRAY),
    ("const { data } = await api.get('/receivings', { params });", WHITE),
    ("", WHITE),
    ("// POST (생성) / PATCH (수정)", LIGHT_GRAY),
    ("const { data } = await api.post('/auth/login', body);", WHITE),
    ("const { data } = await api.patch(`/receivings/${id}/confirm`);", WHITE),
], font_size=10)

# 응답 섹션
add_text(slide, Inches(6.9), Inches(1.1), Inches(6), Inches(0.4),
         "📥  백 → 프론트  응답 (Response)", font_size=20, bold=True, color=ACCENT_GREEN)

add_shape(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3), fill_color=RGBColor(0x15, 0x2B, 0x15), border_color=ACCENT_GREEN, border_width=Pt(1))

# 성공 응답
add_text(slide, Inches(7.1), Inches(1.75), Inches(5.4), Inches(0.35), "✅  성공 응답 (200)", font_size=14, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(7.1), Inches(2.15), Inches(5.4), Inches(1.2), [
    ('{', ACCENT_GREEN),
    ('  "success": true,', WHITE),
    ('  "data": {', WHITE),
    ('    "accessToken": "eyJhbG...",', ACCENT_YELLOW),
    ('    "refreshToken": "eyJhbG..."', ACCENT_YELLOW),
    ('  }', WHITE),
    ('}', ACCENT_GREEN),
], font_size=11)

# 목록 응답
add_text(slide, Inches(7.1), Inches(3.5), Inches(5.4), Inches(0.35), "📋  목록 응답 (페이지네이션)", font_size=14, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(7.1), Inches(3.9), Inches(5.4), Inches(1.5), [
    ('{', ACCENT_BLUE),
    ('  "success": true,', WHITE),
    ('  "data": [ {...}, {...}, ... ],', ACCENT_YELLOW),
    ('  "meta": {', WHITE),
    ('    "page": 1, "limit": 20,', LIGHT_GRAY),
    ('    "total": 45, "totalPages": 3', LIGHT_GRAY),
    ('  }', WHITE),
    ('}', ACCENT_BLUE),
], font_size=11)

# 에러 응답
add_text(slide, Inches(7.1), Inches(5.55), Inches(5.4), Inches(0.35), "❌  에러 응답 (4xx / 5xx)", font_size=14, bold=True, color=ACCENT_RED)
add_code_block(slide, Inches(7.1), Inches(5.95), Inches(5.4), Inches(0.8), [
    ('{', ACCENT_RED),
    ('  "success": false,', WHITE),
    ('  "error": { "code": "INVALID_CREDENTIALS",', ACCENT_RED),
    ('             "message": "이메일 또는 비밀번호 오류" }', ACCENT_RED),
    ('}', ACCENT_RED),
], font_size=11)


# ════════════════════════════════════════
# SLIDE 4: 로그인 통신 예시 (상세)
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "예시 ①  로그인 — 전체 흐름 + 코드", font_size=28, bold=True, color=WHITE)
add_pill(slide, Inches(8.5), Inches(0.35), Inches(1.6), Inches(0.38), ACCENT_BLUE, "POST", font_size=13)
add_text(slide, Inches(10.1), Inches(0.35), Inches(3), Inches(0.38), "/api/auth/login", font_size=16, bold=True, color=VERY_LIGHT, font_name="Consolas")

# 흐름도 (수평 스텝)
steps = [
    ("1", "프론트\nlogin/page.tsx", ACCENT_BLUE, "api.post(\n '/auth/login',\n { email, password }\n)"),
    ("2", "Next.js\nnext.config.js", ACCENT_PURPLE, "rewrite\n/api/* → :4003"),
    ("3", "Gateway\nmain.ts", ACCENT_ORANGE, "/api/auth\n→ :4001 프록시"),
    ("4", "Auth API\nauth.controller.ts", ACCENT_GREEN, "@Post('login')\nvalidate → bcrypt\n→ JWT 발급"),
]

for i, (num, label, clr, desc) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    # 원형 번호
    add_circle(slide, x + Inches(0.9), Inches(1.2), Inches(0.45), clr, num, font_size=18)
    # 레이블
    add_text(slide, x, Inches(1.75), Inches(2.8), Inches(0.5), label, font_size=12, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    # 설명
    add_text(slide, x, Inches(2.25), Inches(2.8), Inches(0.7), desc, font_size=10, color=VERY_LIGHT, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    # 화살표
    if i < 3:
        add_text(slide, x + Inches(2.5), Inches(1.25), Inches(0.7), Inches(0.4), "→", font_size=24, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 프론트 코드
add_text(slide, Inches(0.6), Inches(3.2), Inches(6), Inches(0.4),
         "프론트엔드 코드  📂 app/login/page.tsx", font_size=16, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(0.6), Inches(3.65), Inches(6.0), Inches(3.3), [
    ("const handleLogin = async () => {", WHITE),
    ("  try {", WHITE),
    ("    // 1. 로그인 요청", LIGHT_GRAY),
    ("    const { data } = await api.post('/auth/login', {", ACCENT_YELLOW),
    ("      email, password", ACCENT_YELLOW),
    ("    });", ACCENT_YELLOW),
    ("", WHITE),
    ("    // 2. 토큰 저장 (Zustand)", LIGHT_GRAY),
    ("    setTokens(data.data.accessToken,", ACCENT_GREEN),
    ("             data.data.refreshToken);", ACCENT_GREEN),
    ("", WHITE),
    ("    // 3. 프로필 조회 (JWT 자동 첨부)", LIGHT_GRAY),
    ("    const profile = await api.get('/auth/profile');", ACCENT_BLUE),
    ("    setUser(profile.data.data);", ACCENT_BLUE),
    ("", WHITE),
    ("    // 4. 역할별 리다이렉트", LIGHT_GRAY),
    ("    router.push(role === 'bank' ? '/bank/dashboard'", WHITE),
    ("                : role === 'mart' ? '/mart/receiving'", WHITE),
    ("                : '/admin/dashboard');", WHITE),
    ("  } catch (err) { setError('로그인 실패'); }", ACCENT_RED),
    ("};", WHITE),
], font_size=10)

# 백엔드 코드
add_text(slide, Inches(6.9), Inches(3.2), Inches(6), Inches(0.4),
         "백엔드 코드  📂 auth.controller.ts + service.ts", font_size=16, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(6.9), Inches(3.65), Inches(6.0), Inches(3.3), [
    ("// Controller", LIGHT_GRAY),
    ("@Post('login')", ACCENT_PURPLE),
    ("async login(@Body() dto: LoginDto) {", WHITE),
    ("  const tokens = await this.authService.login(dto);", WHITE),
    ("  return createResponse(tokens);", ACCENT_GREEN),
    ("}  // → { success: true, data: { accessToken, refreshToken } }", LIGHT_GRAY),
    ("", WHITE),
    ("// Service", LIGHT_GRAY),
    ("async login(dto: LoginDto) {", WHITE),
    ("  // 1. 사용자 조회", LIGHT_GRAY),
    ("  const user = await this.userRepo.findOne(", WHITE),
    ("    { where: { email: dto.email } }", WHITE),
    ("  );", WHITE),
    ("  // 2. 비밀번호 검증", LIGHT_GRAY),
    ("  const ok = await bcrypt.compare(dto.password,", ACCENT_YELLOW),
    ("                                  user.passwordHash);", ACCENT_YELLOW),
    ("  // 3. JWT 발급", LIGHT_GRAY),
    ("  const accessToken = this.jwtService.sign(payload);", ACCENT_BLUE),
    ("  const refreshToken = this.jwtService.sign(", ACCENT_BLUE),
    ("    payload, { expiresIn: '7d' });", ACCENT_BLUE),
    ("  return { accessToken, refreshToken };", ACCENT_GREEN),
    ("}", WHITE),
], font_size=10)


# ════════════════════════════════════════
# SLIDE 5: 입고 조회 + 입고확인 예시
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "예시 ②  입고 조회 & ⭐ 입고확인", font_size=28, bold=True, color=WHITE)

# ── 좌측: 입고 조회 ──
add_text(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(0.4),
         "📋  입고 목록 조회", font_size=18, bold=True, color=ACCENT_BLUE)
add_pill(slide, Inches(4.0), Inches(1.05), Inches(1.0), Inches(0.33), ACCENT_GREEN, "GET", font_size=12)
add_text(slide, Inches(5.0), Inches(1.05), Inches(3), Inches(0.33), "/api/receivings", font_size=13, bold=True, color=VERY_LIGHT, font_name="Consolas")

# 프론트 코드
add_text(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.3), "프론트  📂 hooks/use-receivings.ts", font_size=12, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(0.6), Inches(1.85), Inches(6.0), Inches(1.1), [
    ("// React Query로 자동 캐싱 + 자동 갱신", LIGHT_GRAY),
    ("const { data } = useQuery({", WHITE),
    ("  queryKey: ['receivings', filters],", ACCENT_YELLOW),
    ("  queryFn: () => api.get('/receivings', { params: filters })", WHITE),
    ("});", WHITE),
    ("// data.data.data = 입고 목록 배열", LIGHT_GRAY),
    ("// data.data.meta = { page, limit, total, totalPages }", LIGHT_GRAY),
], font_size=10)

# 백엔드 코드
add_text(slide, Inches(0.6), Inches(3.1), Inches(6), Inches(0.3), "백엔드  📂 receiving.controller.ts", font_size=12, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(0.6), Inches(3.45), Inches(6.0), Inches(1.0), [
    ("@Get()", ACCENT_PURPLE),
    ("@UseGuards(JwtAuthGuard)", ACCENT_PURPLE),
    ("async findAll(@Query() query, @CurrentUser() user) {", WHITE),
    ("  const result = await this.service.findAll(query, user);", WHITE),
    ("  return createPaginatedResponse(", ACCENT_GREEN),
    ("    result.items, result.total, query.page, query.limit", ACCENT_GREEN),
    ("  );", ACCENT_GREEN),
    ("}", WHITE),
], font_size=10)

# ── 우측: 입고확인 ──
add_text(slide, Inches(6.9), Inches(1.0), Inches(6), Inches(0.4),
         "⭐  입고확인 (핵심 트리거!)", font_size=18, bold=True, color=ACCENT_YELLOW)
add_pill(slide, Inches(10.5), Inches(1.05), Inches(1.3), Inches(0.33), ACCENT_ORANGE, "PATCH", font_size=12)
add_text(slide, Inches(11.8), Inches(1.05), Inches(1.5), Inches(0.33), "/:id/confirm", font_size=13, bold=True, color=VERY_LIGHT, font_name="Consolas")

# 프론트 코드
add_text(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(0.3), "프론트  📂 hooks/use-receivings.ts", font_size=12, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(6.9), Inches(1.85), Inches(6.0), Inches(1.1), [
    ("// 입고확인 버튼 클릭 시", LIGHT_GRAY),
    ("const confirmMutation = useMutation({", WHITE),
    ("  mutationFn: (id: string) =>", WHITE),
    ("    api.patch(`/receivings/${id}/confirm`),", ACCENT_YELLOW),
    ("  onSuccess: () => {", WHITE),
    ("    queryClient.invalidateQueries(['receivings']);", ACCENT_GREEN),
    ("  }", WHITE),
    ("});", WHITE),
], font_size=10)

# 백엔드 코드
add_text(slide, Inches(6.9), Inches(3.1), Inches(6), Inches(0.3), "백엔드  📂 receiving.service.ts  (핵심 비즈니스 로직)", font_size=12, bold=True, color=ACCENT_YELLOW)
add_code_block(slide, Inches(6.9), Inches(3.45), Inches(6.0), Inches(1.0), [
    ("@Patch(':id/confirm')  @UseGuards(JwtAuthGuard)", ACCENT_PURPLE),
    ("async confirm(@Param('id') id, @CurrentUser() user) {", WHITE),
    ("  // 비관적 잠금 + 트랜잭션으로 상태 변경", LIGHT_GRAY),
    ("  // → Bull Queue에 신용점수 재산출 작업 추가", LIGHT_GRAY),
    ("  // → WebSocket으로 은행에 실시간 알림", LIGHT_GRAY),
    ("  return createResponse(", ACCENT_GREEN),
    ("    await this.service.confirmReceiving(id, user)", ACCENT_GREEN),
    ("  );", ACCENT_GREEN),
    ("}", WHITE),
], font_size=10)

# 하단: 입고확인 후 일어나는 일
add_text(slide, Inches(0.6), Inches(4.75), Inches(12), Inches(0.4),
         "⭐  입고확인 클릭 후 일어나는 일 (시스템 핵심 흐름)", font_size=18, bold=True, color=ACCENT_YELLOW)

flow_steps = [
    ("1", "입고확인\n버튼 클릭", ACCENT_BLUE),
    ("2", "비관적 잠금\n상태→confirmed", ACCENT_ORANGE),
    ("3", "Bull Queue\n신용점수 재산출", ACCENT_PURPLE),
    ("4", "WebSocket\n은행 실시간 알림", ACCENT_GREEN),
    ("5", "프론트 캐시\n무효화→화면갱신", ACCENT_YELLOW),
]

for i, (num, label, clr) in enumerate(flow_steps):
    x = Inches(0.5 + i * 2.5)
    add_shape(slide, x, Inches(5.25), Inches(2.1), Inches(1.0), fill_color=BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, x + Inches(0.05), Inches(5.1), Inches(0.4), clr, num, font_size=14)
    add_text(slide, x, Inches(5.35), Inches(2.1), Inches(0.8), label, font_size=11, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + Inches(2.0), Inches(5.45), Inches(0.6), Inches(0.4), "→", font_size=22, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 6: WebSocket 통신
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "WebSocket 통신  —  백엔드 → 프론트  실시간 알림", font_size=28, bold=True, color=WHITE)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "백엔드에서 이벤트를 보내면, 프론트는 해당 데이터 캐시만 무효화 → 화면이 자동 갱신됩니다", font_size=14, color=LIGHT_GRAY)

# 흐름도
add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
         "실시간 알림 흐름도", font_size=18, bold=True, color=ACCENT_GREEN)

ws_steps = [
    ("erp-api\n입고확인 완료!", ACCENT_ORANGE, "📂 receiving\n.service.ts"),
    ("NotificationGateway\nnotifyBank(event, data)", ACCENT_GREEN, "📂 notification\n.gateway.ts"),
    ("Socket.IO\n'bank' room 전송", ACCENT_PURPLE, "server.to('bank')\n.emit(event, data)"),
    ("프론트\nsocket.on(event)", ACCENT_BLUE, "📂 hooks/\nuse-socket.ts"),
    ("React Query\n캐시 무효화→갱신", ACCENT_YELLOW, "invalidateQueries\n(['receivings'])"),
]

for i, (label, clr, desc) in enumerate(ws_steps):
    x = Inches(0.3 + i * 2.6)
    add_shape(slide, x, Inches(2.0), Inches(2.3), Inches(1.5), fill_color=BG_CARD, border_color=clr, border_width=Pt(2))
    add_text(slide, x, Inches(2.1), Inches(2.3), Inches(0.7), label, font_size=12, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.85), Inches(2.3), Inches(0.55), desc, font_size=9, color=LIGHT_GRAY, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + Inches(2.1), Inches(2.35), Inches(0.6), Inches(0.4), "→", font_size=22, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 백엔드 코드
add_text(slide, Inches(0.6), Inches(3.75), Inches(6), Inches(0.4),
         "백엔드 — 이벤트 보내기  📂 notification.gateway.ts", font_size=16, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(0.6), Inches(4.2), Inches(6.0), Inches(2.0), [
    ("@WebSocketGateway({ cors: true })", ACCENT_PURPLE),
    ("export class NotificationGateway {", WHITE),
    ("  @WebSocketServer() server: Server;", WHITE),
    ("", WHITE),
    ("  // 연결 시 role 기반으로 room 분리", LIGHT_GRAY),
    ("  handleConnection(client) {", WHITE),
    ("    const role = verifyJwt(token).role;", WHITE),
    ("    client.join(role); // 'bank', 'mart', 'admin'", ACCENT_YELLOW),
    ("  }", WHITE),
    ("", WHITE),
    ("  // 은행에게만 이벤트 전송", LIGHT_GRAY),
    ("  notifyBank(event: string, data: unknown) {", WHITE),
    ("    this.server.to('bank').emit(event, data);", ACCENT_GREEN),
    ("  }", WHITE),
    ("}", WHITE),
], font_size=10)

# 프론트 코드
add_text(slide, Inches(6.9), Inches(3.75), Inches(6), Inches(0.4),
         "프론트 — 이벤트 듣기  📂 hooks/use-socket.ts", font_size=16, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(6.9), Inches(4.2), Inches(6.0), Inches(2.0), [
    ("import { getSocket } from '@/lib/socket';", ACCENT_PURPLE),
    ("import { useQueryClient } from '@tanstack/react-query';", ACCENT_PURPLE),
    ("", WHITE),
    ("export function useSocket() {", WHITE),
    ("  const queryClient = useQueryClient();", WHITE),
    ("  const socket = getSocket();", WHITE),
    ("", WHITE),
    ("  // 이벤트 이름에 맞춰 캐시 무효화", LIGHT_GRAY),
    ("  socket.on('receiving:confirmed', () => {", ACCENT_GREEN),
    ("    queryClient.invalidateQueries(['receivings']);", ACCENT_YELLOW),
    ("    queryClient.invalidateQueries(['dashboard']);", ACCENT_YELLOW),
    ("  });", WHITE),
    ("  socket.on('credit:score-updated', () => {", ACCENT_GREEN),
    ("    queryClient.invalidateQueries(['merchants']);", ACCENT_YELLOW),
    ("  });", WHITE),
    ("}", WHITE),
], font_size=10)

# 이벤트 매핑 테이블
add_text(slide, Inches(0.6), Inches(6.4), Inches(12), Inches(0.35),
         "📌  이벤트 ↔ 무효화 데이터 매핑", font_size=15, bold=True, color=ACCENT_YELLOW)

table_data = [
    ("이벤트 이름", "누가 보냄", "프론트에서 무효화", "화면 갱신"),
    ("receiving:confirmed", "erp-api", "receivings, dashboard", "입고목록, 대시보드"),
    ("credit:score-updated", "engine (via erp)", "merchants", "가맹점 목록"),
    ("risk:level-changed", "engine (via erp)", "marts, branches", "마트, 지점 현황"),
]

# 간단한 테이블 (shape 기반)
for row_i, row in enumerate(table_data):
    for col_i, cell in enumerate(row):
        col_widths = [Inches(2.8), Inches(2.5), Inches(3.5), Inches(3.5)]
        col_x = Inches(0.6) + sum(col_widths[:col_i])
        y = Inches(6.75) + row_i * Inches(0.3)
        clr = ACCENT_BLUE if row_i == 0 else WHITE
        sz = 10 if row_i == 0 else 9
        bld = row_i == 0
        fn = "맑은 고딕" if row_i == 0 else "Consolas"
        add_text(slide, col_x, y, col_widths[col_i], Inches(0.25), cell, font_size=sz, bold=bld, color=clr, font_name=fn)


# ════════════════════════════════════════
# SLIDE 7: JSON 요청/응답 형식 정리
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "JSON 형식 정리  —  요청 & 응답 규격", font_size=28, bold=True, color=WHITE)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "모든 API가 이 형식을 따릅니다. 새 API를 만들 때도 같은 형식으로 만들면 됩니다.", font_size=14, color=LIGHT_GRAY)

# 좌측: 요청 형식들
add_text(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
         "📤  요청 형식 (프론트 → 백)", font_size=18, bold=True, color=ACCENT_BLUE)

# GET 요청
add_text(slide, Inches(0.6), Inches(1.9), Inches(6), Inches(0.3), "GET 요청 — URL query로 조건 전달", font_size=13, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(0.6), Inches(2.25), Inches(6.0), Inches(0.85), [
    ("GET /api/receivings?martId=uuid&status=pending&page=1&limit=20", ACCENT_YELLOW),
    ("Headers:", LIGHT_GRAY),
    ("  Authorization: Bearer eyJhbGciOiJIUzI1NiIs...", ACCENT_BLUE),
    ("  Content-Type: application/json", WHITE),
    ("Body: 없음", LIGHT_GRAY),
], font_size=10)

# POST 요청
add_text(slide, Inches(0.6), Inches(3.25), Inches(6), Inches(0.3), "POST 요청 — Body에 JSON 데이터", font_size=13, bold=True, color=ACCENT_ORANGE)
add_code_block(slide, Inches(0.6), Inches(3.6), Inches(6.0), Inches(1.1), [
    ("POST /api/auth/login", ACCENT_YELLOW),
    ("Headers:", LIGHT_GRAY),
    ("  Content-Type: application/json", WHITE),
    ("Body:", LIGHT_GRAY),
    ('{', WHITE),
    ('  "email": "admin@itadx.com",', ACCENT_GREEN),
    ('  "password": "password123"', ACCENT_GREEN),
    ('}', WHITE),
], font_size=10)

# PATCH 요청
add_text(slide, Inches(0.6), Inches(4.85), Inches(6), Inches(0.3), "PATCH 요청 — URL에 ID, Body 선택적", font_size=13, bold=True, color=ACCENT_PURPLE)
add_code_block(slide, Inches(0.6), Inches(5.2), Inches(6.0), Inches(0.65), [
    ("PATCH /api/receivings/uuid-123/confirm", ACCENT_YELLOW),
    ("Headers:", LIGHT_GRAY),
    ("  Authorization: Bearer eyJhbGciOiJIUzI1NiIs...", ACCENT_BLUE),
    ("Body: 없음 (또는 수정할 필드만)", LIGHT_GRAY),
], font_size=10)

# 우측: 응답 형식들
add_text(slide, Inches(6.9), Inches(1.4), Inches(6), Inches(0.4),
         "📥  응답 형식 (백 → 프론트)  — 3가지만 있음!", font_size=18, bold=True, color=ACCENT_GREEN)

# 단일 객체 응답
add_text(slide, Inches(6.9), Inches(1.9), Inches(6), Inches(0.3), "① 단일 객체 응답", font_size=13, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(6.9), Inches(2.25), Inches(6.0), Inches(0.75), [
    ('{ "success": true,', WHITE),
    ('  "data": {                    ← 항상 여기에 결과', ACCENT_YELLOW),
    ('    "accessToken": "...", "refreshToken": "..."', ACCENT_GREEN),
    ('  }', WHITE),
    ('}', WHITE),
], font_size=10)

# 목록 + 페이지네이션 응답
add_text(slide, Inches(6.9), Inches(3.15), Inches(6), Inches(0.3), "② 목록 + 페이지네이션 응답", font_size=13, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(6.9), Inches(3.5), Inches(6.0), Inches(1.1), [
    ('{ "success": true,', WHITE),
    ('  "data": [                    ← 배열', ACCENT_YELLOW),
    ('    { "id": "...", "status": "pending", ... },', ACCENT_GREEN),
    ('    { "id": "...", "status": "confirmed", ... }', ACCENT_GREEN),
    ('  ],', WHITE),
    ('  "meta": {                    ← 페이지 정보', ACCENT_BLUE),
    ('    "page": 1, "limit": 20, "total": 45, "totalPages": 3', ACCENT_BLUE),
    ('  }', WHITE),
    ('}', WHITE),
], font_size=10)

# 에러 응답
add_text(slide, Inches(6.9), Inches(4.75), Inches(6), Inches(0.3), "③ 에러 응답", font_size=13, bold=True, color=ACCENT_RED)
add_code_block(slide, Inches(6.9), Inches(5.1), Inches(6.0), Inches(0.75), [
    ('{ "success": false,', WHITE),
    ('  "error": {', ACCENT_RED),
    ('    "code": "INVALID_CREDENTIALS",    ← 에러 코드', ACCENT_RED),
    ('    "message": "비밀번호가 올바르지 않습니다"', ACCENT_RED),
    ('  }', ACCENT_RED),
    ('}', WHITE),
], font_size=10)

# 하단: 백엔드 공통 유틸
add_text(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.35),
         "💡  백엔드: createResponse(data) / createPaginatedResponse(items, total, page, limit)  쓰면 위 형식이 자동 적용됩니다", font_size=13, bold=True, color=ACCENT_YELLOW)
add_text(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.35),
         "💡  프론트: data.data (단일/목록) / data.meta (페이지) / data.error (에러) 로 접근하면 됩니다", font_size=13, bold=True, color=ACCENT_BLUE)


# ════════════════════════════════════════
# SLIDE 8: 인증(JWT) 흐름
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "JWT 인증 흐름  —  토큰 자동 관리", font_size=28, bold=True, color=WHITE)

# 타임라인 흐름
jwt_steps = [
    ("1", "로그인", "POST /auth/login\n→ accessToken (15분)\n→ refreshToken (7일)", ACCENT_BLUE),
    ("2", "저장", "Zustand persist\nlocalStorage 저장\n→ 새로고침해도 유지", ACCENT_GREEN),
    ("3", "요청마다", "api interceptor가\nAuthorization: Bearer ...\n헤더 자동 첨부", ACCENT_PURPLE),
    ("4", "만료 시", "401 받으면 자동으로\nPOST /auth/refresh\n→ 새 accessToken 발급", ACCENT_ORANGE),
    ("5", "갱신 실패", "refreshToken도 만료\n→ 로그아웃 처리\n→ /login 리다이렉트", ACCENT_RED),
]

for i, (num, title, desc, clr) in enumerate(jwt_steps):
    x = Inches(0.3 + i * 2.6)
    add_shape(slide, x, Inches(1.2), Inches(2.3), Inches(2.0), fill_color=BG_CARD, border_color=clr, border_width=Pt(2))
    add_circle(slide, x + Inches(0.05), Inches(1.05), Inches(0.4), clr, num, font_size=14)
    add_text(slide, x, Inches(1.35), Inches(2.3), Inches(0.35), title, font_size=14, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(1.75), Inches(2.3), Inches(1.2), desc, font_size=10, color=VERY_LIGHT, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + Inches(2.1), Inches(1.8), Inches(0.6), Inches(0.4), "→", font_size=22, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 코드: api.ts interceptor
add_text(slide, Inches(0.6), Inches(3.5), Inches(6), Inches(0.4),
         "프론트 — JWT 자동 관리 코드  📂 lib/api.ts", font_size=16, bold=True, color=ACCENT_BLUE)
add_code_block(slide, Inches(0.6), Inches(3.95), Inches(6.0), Inches(2.8), [
    ("// 요청 인터셉터: 매 요청에 JWT 자동 첨부", LIGHT_GRAY),
    ("api.interceptors.request.use((config) => {", WHITE),
    ("  const token = useAuthStore.getState().accessToken;", WHITE),
    ("  if (token)", WHITE),
    ("    config.headers.Authorization = `Bearer ${token}`;", ACCENT_YELLOW),
    ("  return config;", WHITE),
    ("});", WHITE),
    ("", WHITE),
    ("// 응답 인터셉터: 401 시 자동 토큰 갱신", LIGHT_GRAY),
    ("api.interceptors.response.use(res => res, async (err) => {", WHITE),
    ("  if (err.response?.status === 401 && !original._retry) {", WHITE),
    ("    const { data } = await axios.post(", ACCENT_GREEN),
    ("      '/api/auth/refresh', { refreshToken }", ACCENT_GREEN),
    ("    );", ACCENT_GREEN),
    ("    // 새 토큰으로 원래 요청 재시도", LIGHT_GRAY),
    ("    return api(original);", ACCENT_BLUE),
    ("  }", WHITE),
    ("  // 갱신도 실패하면 → 로그아웃", LIGHT_GRAY),
    ("  useAuthStore.getState().logout();", ACCENT_RED),
    ("});", WHITE),
], font_size=10)

# 백엔드 인증
add_text(slide, Inches(6.9), Inches(3.5), Inches(6), Inches(0.4),
         "백엔드 — 인증 보호 적용 방법", font_size=16, bold=True, color=ACCENT_GREEN)
add_code_block(slide, Inches(6.9), Inches(3.95), Inches(6.0), Inches(1.4), [
    ("// 인증이 필요한 API에 데코레이터만 붙이면 됨", LIGHT_GRAY),
    ("@UseGuards(JwtAuthGuard)       ← 이것만 붙이면 인증 체크", ACCENT_YELLOW),
    ("@Get()", ACCENT_PURPLE),
    ("async findAll(@CurrentUser() user: UserPayload) {", WHITE),
    ("  // user.sub  = 사용자 ID", LIGHT_GRAY),
    ("  // user.role = 'bank' | 'mart' | 'admin'", LIGHT_GRAY),
    ("  // user.martId = 마트 ID (mart만 있음)", LIGHT_GRAY),
    ("  return this.service.findAll(user);", ACCENT_GREEN),
    ("}", WHITE),
], font_size=10)

# JWT 구조
add_text(slide, Inches(6.9), Inches(5.55), Inches(6), Inches(0.3), "JWT Payload 구조", font_size=13, bold=True, color=ACCENT_YELLOW)
add_code_block(slide, Inches(6.9), Inches(5.9), Inches(6.0), Inches(0.85), [
    ('{', WHITE),
    ('  "sub": "user-uuid",         // 사용자 ID', ACCENT_GREEN),
    ('  "email": "admin@itadx.com", // 이메일', ACCENT_GREEN),
    ('  "role": "admin",            // bank | mart | admin', ACCENT_YELLOW),
    ('  "martId": null,             // 마트 소속 (mart만)', ACCENT_BLUE),
    ('  "exp": 1772108265           // 만료시각 (15분 후)', LIGHT_GRAY),
    ('}', WHITE),
], font_size=10)


# ════════════════════════════════════════
# SLIDE 9: API 라우트 전체 목록
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "API 라우트 전체 목록  —  Gateway 경로 맵", font_size=28, bold=True, color=WHITE)

# Auth API
add_text(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.35),
         "🔐  Auth API  :4001", font_size=16, bold=True, color=ACCENT_BLUE)
auth_routes = [
    ("POST", "/api/auth/login", "로그인 (토큰 발급)", "❌"),
    ("POST", "/api/auth/refresh", "토큰 갱신", "❌"),
    ("POST", "/api/auth/logout", "로그아웃", "✅"),
    ("GET", "/api/auth/profile", "내 프로필 조회", "✅"),
]
y = Inches(1.5)
# 헤더
for col_i, header in enumerate(["메서드", "URL", "설명", "인증"]):
    col_widths_auth = [Inches(1.0), Inches(2.8), Inches(2.0), Inches(0.6)]
    col_x = Inches(0.6) + sum(col_widths_auth[:col_i])
    add_text(slide, col_x, y, col_widths_auth[col_i], Inches(0.25), header, font_size=10, bold=True, color=ACCENT_BLUE)
y += Inches(0.28)
for method, url, desc, auth in auth_routes:
    mcol = ACCENT_GREEN if method == "GET" else ACCENT_ORANGE if method == "POST" else ACCENT_PURPLE
    for col_i, (cell, clr_c) in enumerate([(method, mcol), (url, VERY_LIGHT), (desc, LIGHT_GRAY), (auth, WHITE)]):
        col_x = Inches(0.6) + sum(col_widths_auth[:col_i])
        fn = "Consolas" if col_i <= 1 else "맑은 고딕"
        add_text(slide, col_x, y, col_widths_auth[col_i], Inches(0.22), cell, font_size=9, color=clr_c, font_name=fn)
    y += Inches(0.23)

# Admin API
add_text(slide, Inches(0.6), Inches(3.0), Inches(6), Inches(0.35),
         "🏢  Admin API  :4000", font_size=16, bold=True, color=ACCENT_GREEN)
admin_routes = [
    ("GET", "/api/marts", "마트 목록"),
    ("GET", "/api/marts/:id", "마트 상세"),
    ("GET", "/api/branches", "지점 목록"),
    ("GET", "/api/merchants", "가맹점 목록"),
    ("GET", "/api/dashboard/summary", "대시보드 요약"),
    ("GET", "/api/users", "사용자 목록"),
]
y = Inches(3.4)
for col_i, header in enumerate(["메서드", "URL", "설명"]):
    col_widths_adm = [Inches(1.0), Inches(3.0), Inches(2.2)]
    col_x = Inches(0.6) + sum(col_widths_adm[:col_i])
    add_text(slide, col_x, y, col_widths_adm[col_i], Inches(0.25), header, font_size=10, bold=True, color=ACCENT_GREEN)
y += Inches(0.28)
for method, url, desc in admin_routes:
    for col_i, (cell, clr_c) in enumerate([(method, ACCENT_GREEN), (url, VERY_LIGHT), (desc, LIGHT_GRAY)]):
        col_x = Inches(0.6) + sum(col_widths_adm[:col_i])
        fn = "Consolas" if col_i <= 1 else "맑은 고딕"
        add_text(slide, col_x, y, col_widths_adm[col_i], Inches(0.22), cell, font_size=9, color=clr_c, font_name=fn)
    y += Inches(0.23)

# ERP API
add_text(slide, Inches(6.9), Inches(1.1), Inches(6), Inches(0.35),
         "⭐  ERP API  :4002  (핵심)", font_size=16, bold=True, color=ACCENT_YELLOW)
erp_routes = [
    ("GET", "/api/receivings", "입고 목록 (필터+페이지)"),
    ("GET", "/api/receivings/:id", "입고 상세"),
    ("PATCH", "/api/receivings/:id/confirm", "⭐ 입고확인"),
    ("GET", "/api/settlements", "정산 목록"),
    ("GET", "/api/settlements/:id", "정산 상세"),
]
y = Inches(1.5)
for col_i, header in enumerate(["메서드", "URL", "설명"]):
    col_widths_erp = [Inches(1.0), Inches(3.2), Inches(2.0)]
    col_x = Inches(6.9) + sum(col_widths_erp[:col_i])
    add_text(slide, col_x, y, col_widths_erp[col_i], Inches(0.25), header, font_size=10, bold=True, color=ACCENT_YELLOW)
y += Inches(0.28)
for method, url, desc in erp_routes:
    mcol = ACCENT_GREEN if method == "GET" else ACCENT_PURPLE
    for col_i, (cell, clr_c) in enumerate([(method, mcol), (url, VERY_LIGHT), (desc, LIGHT_GRAY)]):
        col_x = Inches(6.9) + sum(col_widths_erp[:col_i])
        fn = "Consolas" if col_i <= 1 else "맑은 고딕"
        add_text(slide, col_x, y, col_widths_erp[col_i], Inches(0.22), cell, font_size=9, color=clr_c, font_name=fn)
    y += Inches(0.23)

# Engine API
add_text(slide, Inches(6.9), Inches(3.0), Inches(6), Inches(0.35),
         "🧠  Engine API  :8000  (Python)", font_size=16, bold=True, color=ACCENT_PURPLE)
engine_routes = [
    ("POST", "/api/v41/screen", "v41 심사 실행"),
    ("POST", "/api/v10/score", "v10 신용점수 산출"),
    ("POST", "/api/dual-track/assess", "듀얼트랙 리스크 평가"),
    ("GET", "/api/branch-risk/:martId", "점포 리스크 조회"),
]
y = Inches(3.4)
for col_i, header in enumerate(["메서드", "URL", "설명"]):
    col_widths_eng = [Inches(1.0), Inches(3.2), Inches(2.0)]
    col_x = Inches(6.9) + sum(col_widths_eng[:col_i])
    add_text(slide, col_x, y, col_widths_eng[col_i], Inches(0.25), header, font_size=10, bold=True, color=ACCENT_PURPLE)
y += Inches(0.28)
for method, url, desc in engine_routes:
    mcol = ACCENT_GREEN if method == "GET" else ACCENT_ORANGE
    for col_i, (cell, clr_c) in enumerate([(method, mcol), (url, VERY_LIGHT), (desc, LIGHT_GRAY)]):
        col_x = Inches(6.9) + sum(col_widths_eng[:col_i])
        fn = "Consolas" if col_i <= 1 else "맑은 고딕"
        add_text(slide, col_x, y, col_widths_eng[col_i], Inches(0.22), cell, font_size=9, color=clr_c, font_name=fn)
    y += Inches(0.23)

# 하단 요약
add_shape(slide, Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.5), fill_color=RGBColor(0x1A, 0x2A, 0x1A), border_color=ACCENT_GREEN, border_width=Pt(1))
add_text(slide, Inches(0.8), Inches(5.6), Inches(12), Inches(0.35),
         "💡  Gateway 라우팅 규칙  (gateway-api/main.ts)", font_size=14, bold=True, color=ACCENT_YELLOW)
add_text(slide, Inches(0.8), Inches(5.95), Inches(12), Inches(0.9),
         "/api/auth        →  auth-api   (localhost:4001)    인증 관련\n"
         "/api/marts, /api/branches, /api/merchants, /api/dashboard, /api/users  →  admin-api (localhost:4000)\n"
         "/api/receivings, /api/settlements  →  erp-api   (localhost:4002)    입고/정산 (핵심)\n"
         "/api/v41, /api/v10, /api/dual-track, /api/branch-risk  →  engine-api (localhost:8000)  AI엔진",
         font_size=11, color=VERY_LIGHT, font_name="Consolas")


# ════════════════════════════════════════
# SLIDE 10: 개발자 체크리스트 / 요약
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "개발자 체크리스트  —  이것만 기억하세요!", font_size=28, bold=True, color=WHITE)

# 프론트엔드 개발자
add_shape(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8), fill_color=RGBColor(0x15, 0x25, 0x3B), border_color=ACCENT_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.45), "🖥  프론트엔드 개발자", font_size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

fe_items = [
    ("✅  api.get / api.post / api.patch 만 사용", "→ JWT, 에러처리 다 자동"),
    ("✅  URL은 /api 빼고 경로만", "api.get('/receivings')  api.post('/auth/login')"),
    ("✅  응답은 항상 data.data 로 접근", "성공: data.data  목록: data.meta  에러: data.error"),
    ("✅  React Query로 캐싱 + 자동 갱신", "queryKey 잘 정하면 WebSocket에서 알아서 무효화"),
    ("✅  인증 상태는 Zustand auth-store", "setTokens() → setUser() → router.push()"),
    ("✅  WebSocket은 use-socket.ts만 확인", "이벤트 수신 → invalidateQueries → 화면 자동 갱신"),
]

y = Inches(1.7)
for title, desc in fe_items:
    add_text(slide, Inches(0.7), y, Inches(5.6), Inches(0.3), title, font_size=12, bold=True, color=WHITE)
    add_text(slide, Inches(0.9), y + Inches(0.3), Inches(5.4), Inches(0.25), desc, font_size=10, color=LIGHT_GRAY, font_name="Consolas")
    y += Inches(0.65)

# 백엔드 개발자
add_shape(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8), fill_color=RGBColor(0x15, 0x2B, 0x15), border_color=ACCENT_GREEN, border_width=Pt(2))
add_text(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.45), "🏗  백엔드 개발자", font_size=20, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

be_items = [
    ("✅  응답은 createResponse(data) 사용", "→ { success: true, data: ... } 자동 래핑"),
    ("✅  목록은 createPaginatedResponse()", "→ data + meta(page, limit, total, totalPages)"),
    ("✅  에러는 BusinessException + ErrorCodes", "throw new BusinessException(code, message, status)"),
    ("✅  인증은 @UseGuards(JwtAuthGuard)", "→ @CurrentUser() 로 사용자 정보 바로 접근"),
    ("✅  감사로그는 @Auditable 데코레이터", "중요 변경사항 자동 기록"),
    ("✅  WebSocket은 NotificationGateway", "notifyBank / notifyMart / notifyAll 사용"),
]

y = Inches(1.7)
for title, desc in be_items:
    add_text(slide, Inches(7.0), y, Inches(5.6), Inches(0.3), title, font_size=12, bold=True, color=WHITE)
    add_text(slide, Inches(7.2), y + Inches(0.3), Inches(5.4), Inches(0.25), desc, font_size=10, color=LIGHT_GRAY, font_name="Consolas")
    y += Inches(0.65)

# 하단 핵심 파일
add_text(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.3),
         "📂 핵심 파일:  lib/api.ts  |  next.config.js  |  gateway-api/main.ts  |  auth.controller.ts  |  receiving.controller.ts  |  notification.gateway.ts",
         font_size=11, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 저장
# ════════════════════════════════════════
output_path = r"D:\ita-dev\itadx-system-architecture-claudeCli-test\ItaDX_프론트백_통신가이드.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드: {len(prs.slides)}장")
